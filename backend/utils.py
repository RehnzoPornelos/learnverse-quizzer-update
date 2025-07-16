from pptx import Presentation
import fitz  # PyMuPDF
import docx


def extract_text_from_file(file_path: str) -> str:
    """
    Extract text from PDF, DOCX, PPTX, or TXT files based on file extension.
    """
    ext = file_path.lower().split('.')[-1]
    text = ""

    if ext == 'pdf':
        # PDF via PyMuPDF
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
    elif ext in ('docx', 'doc'):
        # DOCX via python-docx
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + '\n'
    elif ext in ('pptx', 'ppt'):
        # PPTX via python-pptx
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + ' '
                    text += '\n'
    elif ext == 'txt':
        # Plain text
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
    else:
        raise ValueError(f"Unsupported file extension: .{ext}")

    return text.strip()
